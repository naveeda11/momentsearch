"""Document acquisition + parsing + chunking — papers (PDF) and decks (PDF/PPTX).

Pure parsing logic, no Prefect imports (mirrors the fetch/frames/transcript
split on the video side). The Prefect flow lives in doc_pipeline.py.

The locator is decided HERE, at parse time: every page record carries a 1-based
`page` (papers) or `slide` (decks) number that rides through chunking into the
Qdrant payload — a citation is stored alongside its embedding, never computed
at query time.

Crash resume: parse results (text + captions) checkpoint to object storage at
parsed/{user}/{doc_id}.json, keyed by source_hash. A re-run after a killed
worker loads the checkpoint instead of re-parsing, and the enrich stage skips
pages that already have captions.
"""
from __future__ import annotations

import io
import json
import urllib.request
from dataclasses import asdict, dataclass, field
from pathlib import Path

from .. import config, storage
from .fetch import scratch_dir, sha256_file

# Content sniffing — the accept path is shape-only (202 instantly), so junk
# URIs (bench posts 30 example.com probes per run) surface HERE as a clean,
# non-retryable parse failure instead of poisoning the queue.
_PDF_MAGIC = b"%PDF"
_ZIP_MAGIC = b"PK\x03\x04"  # pptx (and any OOXML) is a zip container

_DOWNLOAD_TIMEOUT_S = 120
_MAX_DOWNLOAD_MB = 200


@dataclass
class PageRec:
    """One page (paper) or slide (deck)."""
    page: int                    # 1-based locator
    text: str = ""
    caption: str = ""            # vision-LLM caption (enrich stage fills this)
    needs_caption: bool = False  # image-only page/slide


class UnsupportedDocument(Exception):
    """Body is not a PDF/PPTX — permanent, not worth retrying."""


def infer_ext(uri: str, kind: str) -> str:
    suffix = Path(uri.split("?", 1)[0]).suffix.lower()
    if suffix in (".pdf", ".pptx"):
        return suffix
    return ".pdf"  # papers are PDFs; arxiv-style URIs often have no extension


def sniff_kind_ext(head: bytes) -> str | None:
    if head.startswith(_PDF_MAGIC):
        return ".pdf"
    if head.startswith(_ZIP_MAGIC):
        return ".pptx"
    return None


def fetch_document(row: dict, user_id: str, doc_id: str) -> tuple[Path, str, str]:
    """Acquire the document into worker scratch. Returns (path, source_hash,
    storage_key).

    Checkpoint behaviour: the first successful fetch uploads the original to
    docs/{user}/{doc_id}.{ext}; a re-run finds it there and skips the network.
    """
    dest_hint = infer_ext(row.get("url") or row.get("storage_key") or "", row.get("kind") or "paper")

    key = row.get("storage_key")
    if key and storage.exists(key):
        dest = scratch_dir() / f"{doc_id}{Path(key).suffix or dest_hint}"
        path = storage.download_to(key, dest)
        print(f"[fetch] {doc_id}: resume from storage checkpoint {key}")
        return path, sha256_file(path), key

    uri = row.get("url") or ""
    if not uri.startswith(("http://", "https://")):
        raise UnsupportedDocument(f"unfetchable uri: {uri!r}")
    req = urllib.request.Request(uri, headers={"User-Agent": "momentsearch/1.0"})
    dest = scratch_dir() / f"{doc_id}{dest_hint}"
    size = 0
    with urllib.request.urlopen(req, timeout=_DOWNLOAD_TIMEOUT_S) as resp, dest.open("wb") as out:
        while chunk := resp.read(1 << 20):
            size += len(chunk)
            if size > _MAX_DOWNLOAD_MB * 1024 * 1024:
                dest.unlink(missing_ok=True)
                raise UnsupportedDocument(f"document exceeds {_MAX_DOWNLOAD_MB}MB")
            out.write(chunk)

    head = dest.open("rb").read(8)
    real_ext = sniff_kind_ext(head)
    if real_ext is None:
        dest.unlink(missing_ok=True)
        raise UnsupportedDocument("body is not a PDF or PPTX")
    if real_ext != dest.suffix:
        fixed = dest.with_suffix(real_ext)
        dest.rename(fixed)
        dest = fixed

    # Durable copy = the fetch checkpoint for crash resume (and /api/doc serving).
    key = storage.doc_key(user_id, doc_id, dest.suffix)
    storage.upload_file(dest, key,
                        "application/pdf" if dest.suffix == ".pdf" else config.ALLOWED_DOC_TYPES[1])
    return dest, sha256_file(dest), key


# ── Parsing ───────────────────────────────────────────────────────────────────

def parse_pdf(path: Path) -> tuple[list[PageRec], list[bytes]]:
    """Per-page text + JPEG renders. Returns (pages, renders) index-aligned."""
    import pymupdf

    pages: list[PageRec] = []
    renders: list[bytes] = []
    with pymupdf.open(path) as doc:
        n = min(doc.page_count, config.DOC_MAX_PAGES)
        for i in range(n):
            pg = doc[i]
            text = pg.get_text("text").strip()
            zoom = config.PAGE_RENDER_WIDTH / max(pg.rect.width, 1)
            pix = pg.get_pixmap(matrix=pymupdf.Matrix(zoom, zoom))
            renders.append(pix.tobytes("jpeg"))
            pages.append(PageRec(page=i + 1, text=text,
                                 needs_caption=len(text) < config.CAPTION_MIN_TEXT_CHARS))
    return pages, renders


def parse_pptx(path: Path) -> tuple[list[PageRec], list[bytes | None]]:
    """Per-slide text (shapes + speaker notes). python-pptx cannot rasterize a
    slide, so the 'render' for an image-heavy slide is its largest embedded
    picture (caption input + citation thumbnail); text slides get None."""
    from pptx import Presentation

    pages: list[PageRec] = []
    renders: list[bytes | None] = []
    prs = Presentation(str(path))
    for i, slide in enumerate(list(prs.slides)[:config.DOC_MAX_PAGES]):
        parts: list[str] = []
        biggest: bytes | None = None
        biggest_size = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
            if shape.shape_type == 13:  # PICTURE
                try:
                    blob = shape.image.blob
                    if len(blob) > biggest_size:
                        biggest, biggest_size = blob, len(blob)
                except Exception:
                    pass
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")
        text = "\n".join(parts).strip()
        renders.append(_to_jpeg(biggest) if biggest else None)
        pages.append(PageRec(page=i + 1, text=text,
                             needs_caption=len(text) < config.CAPTION_MIN_TEXT_CHARS))
    return pages, renders


def _to_jpeg(blob: bytes) -> bytes | None:
    """Normalize an embedded picture (png/emf/...) to JPEG; None if unreadable."""
    from PIL import Image

    try:
        img = Image.open(io.BytesIO(blob))
        img.thumbnail((config.PAGE_RENDER_WIDTH, config.PAGE_RENDER_WIDTH * 2))
        buf = io.BytesIO()
        img.convert("RGB").save(buf, format="JPEG", quality=80)
        return buf.getvalue()
    except Exception:
        return None


def parse_document(path: Path):
    if path.suffix == ".pptx":
        return parse_pptx(path)
    return parse_pdf(path)


# ── Chunking ──────────────────────────────────────────────────────────────────

def full_text(p: PageRec) -> str:
    """Text + caption combined — what actually gets embedded."""
    parts = [p.text] if p.text else []
    if p.caption:
        parts.append(f"[Figure] {p.caption}")
    return "\n".join(parts).strip()


def chunk_pages(pages: list[PageRec], kind: str) -> list[dict]:
    """Paper: ~DOC_CHUNK_CHARS chunks that never cross a page boundary.
    Deck: one chunk per slide. Every chunk carries its locator."""
    locator = "slide" if kind == "deck" else "page"
    chunks: list[dict] = []
    for p in pages:
        text = full_text(p)
        if not text:
            continue
        if kind == "deck" or len(text) <= config.DOC_CHUNK_CHARS:
            chunks.append({"text": text, locator: p.page})
            continue
        step = max(config.DOC_CHUNK_CHARS - config.DOC_CHUNK_OVERLAP, 200)
        for start in range(0, len(text), step):
            piece = text[start:start + config.DOC_CHUNK_CHARS].strip()
            if piece:
                chunks.append({"text": piece, locator: p.page})
    return chunks


# ── Parse checkpoint (crash resume) ──────────────────────────────────────────

def save_checkpoint(user_id: str, doc_id: str, source_hash: str,
                    pages: list[PageRec]) -> None:
    payload = {"source_hash": source_hash, "pages": [asdict(p) for p in pages]}
    storage.put_bytes(storage.parsed_key(user_id, doc_id),
                      json.dumps(payload).encode(), "application/json")


def load_checkpoint(user_id: str, doc_id: str, source_hash: str) -> list[PageRec] | None:
    """The parsed pages from a prior run of the SAME content, else None."""
    try:
        raw = storage.get_bytes(storage.parsed_key(user_id, doc_id))
        data = json.loads(raw)
    except Exception:
        return None
    if data.get("source_hash") != source_hash:
        return None  # content changed under the same id — reparse
    return [PageRec(**p) for p in data.get("pages", [])]
