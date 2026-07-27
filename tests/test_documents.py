"""PDF parsing on an in-test pymupdf fixture (no binary files committed)."""
import pathlib
import tempfile

import pymupdf

from src import config
from src.ingest.documents import (
    parse_pdf,
    parse_pptx,
    rendered_pptx_path,
    sniff_kind_ext,
)


def _fixture_pdf() -> pathlib.Path:
    doc = pymupdf.open()
    page1 = doc.new_page()
    line = ("Retrieval augmented generation combines a retriever with a "
            "generator so answers stay grounded in sources.")
    for i in range(6):  # separate lines: insert_text does not wrap
        page1.insert_text((72, 72 + i * 16), line, fontsize=11)
    doc.new_page()  # page 2: image-only stand-in (no extractable text)
    path = pathlib.Path(tempfile.mkdtemp()) / "fixture.pdf"
    doc.save(str(path))
    doc.close()
    return path


def test_parse_pdf_pages_and_captions():
    pages, renders = parse_pdf(_fixture_pdf())
    assert [p.page for p in pages] == [1, 2]
    assert "Retrieval augmented" in pages[0].text
    assert not pages[0].needs_caption          # plenty of text
    assert pages[1].needs_caption              # blank page -> caption candidate
    assert len(renders) == 2 and all(r[:3] == b"\xff\xd8\xff" for r in renders)


def test_render_width():
    pages, renders = parse_pdf(_fixture_pdf())
    import io
    from PIL import Image
    w, _ = Image.open(io.BytesIO(renders[0])).size
    assert abs(w - config.PAGE_RENDER_WIDTH) <= 2


def test_sniff():
    assert sniff_kind_ext(b"%PDF-1.7 ...") == ".pdf"
    assert sniff_kind_ext(b"PK\x03\x04rest") == ".pptx"
    assert sniff_kind_ext(b"<html>") is None


def test_parse_pptx_renders_complete_slides():
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(2))
    shape.text = "Shape-based self-attention diagram"
    shape.fill.solid()
    path = pathlib.Path(tempfile.mkdtemp()) / "fixture.pptx"
    prs.save(path)

    pages, renders = parse_pptx(path)

    assert len(pages) == 1
    assert "self-attention" in pages[0].text
    assert len(renders) == 1 and renders[0][:3] == b"\xff\xd8\xff"
    assert rendered_pptx_path(path).exists()
